{
 "cells": [
  {
   "cell_type": "code",
   "execution_count": null,
   "id": "9db09c99",
   "metadata": {},
   "outputs": [],
   "source": [
    "import openpyxl\n",
    "from pptx import Presentation\n",
    "\n",
    "# Load the Excel file\n",
    "wb = openpyxl.load_workbook(\"C:\\path\\to\\your\\excel\\file.xlsx\")\n",
    "sheet = wb['Sheet1']\n",
    "\n",
    "# Load the PowerPoint template\n",
    "prs = Presentation(\"C:\\path\\to\\your\\powerpoint\\template.pptx\")\n",
    "\n",
    "for i in range(1, 101):\n",
    "    # Update the title and description\n",
    "    title = sheet.cell(row=i, column=1).value\n",
    "    description = sheet.cell(row=i, column=2).value\n",
    "    prs.slides[0].shapes[0].text = title\n",
    "    prs.slides[0].shapes[1].text = description\n",
    "\n",
    "    # Save the slide as a PNG file\n",
    "    filename = \"C:\\path\\to\\save\\images\\thumbnail_{}.png\".format(i)\n",
    "    prs.slides[0].shapes[2\n"
   ]
  }
 ],
 "metadata": {
  "kernelspec": {
   "display_name": "Python 3 (ipykernel)",
   "language": "python",
   "name": "python3"
  },
  "language_info": {
   "codemirror_mode": {
    "name": "ipython",
    "version": 3
   },
   "file_extension": ".py",
   "mimetype": "text/x-python",
   "name": "python",
   "nbconvert_exporter": "python",
   "pygments_lexer": "ipython3",
   "version": "3.9.13"
  }
 },
 "nbformat": 4,
 "nbformat_minor": 5
}
